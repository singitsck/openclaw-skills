#!/usr/bin/env python3
"""
demo_parallel_slides.py - 並行幻燈片創建演示

這個腳本示範如何使用 sessions_spawn 批量創建 PPT 幻燈片
"""

import os
import json
from datetime import datetime

# 模擬課程內容
LESSON_SLIDES = [
    {
        "num": 1,
        "type": "title",
        "title": "什麼是 AI Agent？",
        "subtitle": "給小朋友的人工智能小課堂"
    },
    {
        "num": 2,
        "type": "content",
        "title": "首先，什麼是 AI？",
        "bullets": [
            "AI = Artificial Intelligence（人工智能）",
            "就像電腦裡住了一個聰明的小助手",
            "它可以聽懂你的話，幫你解答問題",
            "例子：Siri、Alexa、手機裡的智慧助理"
        ]
    },
    {
        "num": 3,
        "type": "content",
        "title": "AI 和 Agent 有什麼不同？",
        "bullets": [
            "普通 AI：你問什麼，它答什麼",
            "AI Agent：會主動幫你做很多事情！",
            "就像一個聰明的小管家",
            "它可以自己思考、計劃、執行任務"
        ]
    },
    {
        "num": 4,
        "type": "image",
        "title": "Agent 的超能力",
        "description": "目標導向、拆解問題、使用工具、自主執行"
    }
]

def generate_slide_task(slide_info, output_dir):
    """生成單個幻燈片的創建任務"""
    
    slide_type = slide_info["type"]
    slide_num = slide_info["num"]
    
    if slide_type == "title":
        task = f"""
使用 python-pptx 創建一個標題頁幻燈片：

幻燈片 {slide_num} - 標題頁
- 主標題："{slide_info['title']}"
- 副標題："{slide_info['subtitle']}"
- 背景使用深藍色（RGB: 30, 39, 97）
- 主標題字體：44pt，白色，置中
- 副標題字體：28pt，淺色，置中

保存到：{output_dir}/slide_{slide_num:02d}.pptx

輸出執行結果確認。
"""
    
    elif slide_type == "content":
        bullets_text = "\\n".join([f'- "{b}"' for b in slide_info["bullets"]])
        task = f"""
使用 python-pptx 創建一個內容頁幻燈片：

幻燈片 {slide_num} - 內容頁
- 標題："{slide_info['title']}"
- 內容項目：
{bullets_text}
- 背景：淺藍白（RGB: 240, 248, 255）
- 頂部有藍色標題條（RGB: 2, 128, 144）
- 標題字體：36pt，白色，粗體
- 內容字體：24pt，深灰色

保存到：{output_dir}/slide_{slide_num:02d}.pptx

輸出執行結果確認。
"""
    
    elif slide_type == "image":
        task = f"""
使用 python-pptx 創建一個圖文頁幻燈片：

幻燈片 {slide_num} - 圖文頁
- 大標題："{slide_info['title']}"
- 描述文字："{slide_info['description']}"
- 背景：青色（RGB: 0, 168, 150）
- 中間有一個大白色圓形（模擬圖標區域）
- 標題在圓形下方，40pt，白色，置中
- 描述在最下方，22pt，淺色，置中

保存到：{output_dir}/slide_{slide_num:02d}.pptx

輸出執行結果確認。
"""
    
    return task.strip()

def create_parallel_tasks_config():
    """創建並行任務配置文件"""
    
    output_dir = os.path.expanduser("~/.openclaw/workspace-rem/pptx_parallel_demo")
    os.makedirs(output_dir, exist_ok=True)
    
    tasks = []
    
    print("🚀 生成並行 PPT 創建任務\\n")
    print("=" * 60)
    
    for slide in LESSON_SLIDES:
        task_content = generate_slide_task(slide, output_dir)
        
        task_config = {
            "task_id": f"slide_{slide['num']:02d}",
            "slide_num": slide["num"],
            "slide_type": slide["type"],
            "session_key": f"pptx-creator-slide-{slide['num']:02d}",
            "task": task_content,
            "output_file": f"{output_dir}/slide_{slide['num']:02d}.pptx",
            "created_at": datetime.now().isoformat()
        }
        
        tasks.append(task_config)
        
        print(f"\\n📊 幻燈片 {slide['num']}: {slide['title']}")
        print(f"   類型: {slide['type']}")
        print(f"   Session: {task_config['session_key']}")
    
    # 保存任務配置
    config_file = f"{output_dir}/parallel_tasks_config.json"
    with open(config_file, 'w', encoding='utf-8') as f:
        json.dump({
            "tasks": tasks,
            "total_slides": len(tasks),
            "output_directory": output_dir,
            "created_at": datetime.now().isoformat()
        }, f, indent=2, ensure_ascii=False)
    
    print("\\n" + "=" * 60)
    print(f"\\n✅ 任務配置已生成：{config_file}")
    print(f"📁 輸出目錄：{output_dir}")
    print(f"\\n📋 下一步操作：")
    print("1. 讀取配置文件中的每個任務")
    print("2. 使用 sessions_spawn 啟動每個任務")
    print("3. 等待所有 session 完成")
    print("4. 聚合結果到最終 PPT")
    
    return config_file, tasks

def generate_aggregation_script(output_dir):
    """生成結果聚合腳本"""
    
    script_content = f'''#!/usr/bin/env python3
"""
聚合並行創建的幻燈片到最終 PPT
"""

from pptx import Presentation
from pptx.util import Inches
import os
import glob

OUTPUT_DIR = "{output_dir}"
FINAL_PPT = os.path.join(OUTPUT_DIR, "final_lesson.pptx")

def aggregate_slides():
    """聚合所有幻燈片"""
    
    # 創建新 PPT
    final_prs = Presentation()
    
    # 獲取所有幻燈片文件（按順序）
    slide_files = sorted(glob.glob(os.path.join(OUTPUT_DIR, "slide_*.pptx")))
    
    print(f"🔍 找到 {{len(slide_files)}} 個幻燈片文件")
    
    for slide_file in slide_files:
        print(f"  📄 處理：{{os.path.basename(slide_file)}}")
        
        # 讀取單個幻燈片
        prs = Presentation(slide_file)
        
        # 複製幻燈片到最終 PPT
        for slide in prs.slides:
            # 這裡需要實現幻燈片複製邏輯
            # 簡化版本：僅複製內容文本
            slide_layout = final_prs.slide_layouts[6]  # 空白佈局
            new_slide = final_prs.slides.add_slide(slide_layout)
            
            # 複製文本內容
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    new_shape = new_slide.shapes.add_textbox(
                        shape.left, shape.top, shape.width, shape.height
                    )
                    new_shape.text = shape.text
    
    # 保存最終 PPT
    final_prs.save(FINAL_PPT)
    print(f"\\n✅ 聚合完成：{{FINAL_PPT}}")

if __name__ == "__main__":
    aggregate_slides()
'''
    
    script_file = f"{output_dir}/aggregate_slides.py"
    with open(script_file, 'w') as f:
        f.write(script_content)
    
    os.chmod(script_file, 0o755)
    print(f"\\n📝 聚合腳本已生成：{script_file}")
    
    return script_file

def main():
    """主程序"""
    
    print("=" * 60)
    print("  並行 PPT 創建演示 - Sessions Spawn 方案")
    print("=" * 60)
    
    # 創建並行任務配置
    config_file, tasks = create_parallel_tasks_config()
    
    # 生成聚合腳本
    output_dir = os.path.expanduser("~/.openclaw/workspace-rem/pptx_parallel_demo")
    aggregate_script = generate_aggregation_script(output_dir)
    
    print("\\n" + "=" * 60)
    print("  使用示例代碼")
    print("=" * 60)
    print("""
# 在 OpenClaw 中使用 sessions_spawn 啟動任務：

# 1. 加載任務配置
import json
with open('""" + config_file + """') as f:
    config = json.load(f)

# 2. 並行啟動所有任務
for task in config['tasks']:
    sessions_spawn(
        task=task['task'],
        mode='run',
        session_key=task['session_key']
    )

# 3. 等待完成後聚合結果
exec(command='python """ + aggregate_script + """')
""")
    
    print("\\n✨ 演示配置完成！")

if __name__ == "__main__":
    main()

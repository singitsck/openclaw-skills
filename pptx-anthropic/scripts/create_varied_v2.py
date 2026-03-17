#!/usr/bin/env python3
"""
AI 使用形態演進 - 多樣化版面 FYP PPT (簡化版)
使用 8 種版面設計
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

FONT = "Microsoft JhengHei"

def set_font(run, size_pt, bold=False, color=(15,23,42)):
    run.font.name = FONT
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if isinstance(color, tuple):
        run.font.color.rgb = RGBColor(*color)

def add_title_slide(slide, title, subtitle):
    """版面1: 全屏深色標題頁"""
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 23, 42)
    bg.line.fill.background()
    
    # 裝飾線
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.3), Inches(4), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(59, 130, 246)
    line.line.fill.background()
    
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.3), Inches(12), Inches(1.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_font(p.runs[0], 52, True, (255,255,255))
    
    box = slide.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(12), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = subtitle
    set_font(p.runs[0], 28, False, (148,163,184))

def add_two_column_slide(slide, emoji, title, bullets):
    """版面2: 兩欄圖文"""
    # 背景
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(248, 250, 252)
    bg.line.fill.background()
    
    # 左欄文字
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(6), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = f"{emoji} {title}"
    set_font(p.runs[0], 36, True)
    
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(5))
    tf = box.text_frame
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        set_font(p.runs[0], 22)
        p.space_before = Pt(20)
    
    # 右欄圓形
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(2), Inches(4), Inches(4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(37, 99, 235)
    circle.line.fill.background()
    
    box = slide.shapes.add_textbox(Inches(8), Inches(3.5), Inches(4), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = emoji
    set_font(p.runs[0], 72, False, (255,255,255))
    p.alignment = PP_ALIGN.CENTER

def add_three_cards_slide(slide, title, cards):
    """版面3: 三欄卡片"""
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(241, 245, 249)
    bg.line.fill.background()
    
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(10), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_font(p.runs[0], 36, True)
    
    colors = [(220,38,38), (217,119,6), (25,135,84)]
    for i, (card_title, desc) in enumerate(cards):
        x = Inches(0.6 + i * 4.2)
        
        card = slide.shapes.add_shape(1, x, Inches(2), Inches(3.8), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255,255,255)
        card.line.fill.background()
        
        bar = slide.shapes.add_shape(1, x, Inches(2), Inches(3.8), Pt(6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*colors[i])
        bar.line.fill.background()
        
        box = slide.shapes.add_textbox(x, Inches(2.8), Inches(3.8), Inches(1))
        p = box.text_frame.paragraphs[0]
        p.text = f"0{i+1}"
        set_font(p.runs[0], 48, True, colors[i])
        p.alignment = PP_ALIGN.CENTER
        
        box = slide.shapes.add_textbox(x+Inches(0.2), Inches(3.8), Inches(3.4), Inches(0.8))
        p = box.text_frame.paragraphs[0]
        p.text = card_title
        set_font(p.runs[0], 28, True)
        p.alignment = PP_ALIGN.CENTER
        
        box = slide.shapes.add_textbox(x+Inches(0.2), Inches(4.7), Inches(3.4), Inches(1.5))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        set_font(p.runs[0], 20)
        p.alignment = PP_ALIGN.CENTER

def add_big_circle_slide(slide, emoji, title, points):
    """版面4: 大圓+清單"""
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(240, 253, 244)
    bg.line.fill.background()
    
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(1.5), Inches(4.5), Inches(4.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(25, 135, 84)
    circle.line.fill.background()
    
    box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(4.5), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = emoji
    set_font(p.runs[0], 72, False, (255,255,255))
    p.alignment = PP_ALIGN.CENTER
    
    box = slide.shapes.add_textbox(Inches(5.5), Inches(0.5), Inches(7), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_font(p.runs[0], 40, True, (25,135,84))
    
    box = slide.shapes.add_textbox(Inches(5.5), Inches(1.6), Inches(7), Inches(5.5))
    tf = box.text_frame
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✓ {point}"
        set_font(p.runs[0], 22)
        p.space_before = Pt(24)

def add_grid_slide(slide, title, items):
    """版面5: 2x2 網格"""
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255,255,255)
    bg.line.fill.background()
    
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_font(p.runs[0], 36, True)
    
    colors = [(37,99,235), (111,66,193), (220,38,38), (217,119,6)]
    positions = [(0.5, 1.6), (6.8, 1.6), (0.5, 4.6), (6.8, 4.6)]
    
    for (x, y), (emoji, item_title, desc), color in zip(positions, items, colors):
        rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(5.8), Inches(2.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(*color)
        rect.line.fill.background()
        
        box = slide.shapes.add_textbox(Inches(x), Inches(y+0.3), Inches(5.8), Inches(0.8))
        p = box.text_frame.paragraphs[0]
        p.text = emoji
        set_font(p.runs[0], 40, False, (255,255,255))
        
        box = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+1.1), Inches(5.4), Inches(0.7))
        p = box.text_frame.paragraphs[0]
        p.text = item_title
        set_font(p.runs[0], 28, True, (255,255,255))
        
        box = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+1.8), Inches(5.4), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.text = desc
        set_font(p.runs[0], 16, False, (230,230,230))

def add_sidebar_slide(slide, title, subtitle, points):
    """版面6: 左側大色條"""
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255,255,255)
    bg.line.fill.background()
    
    bar = slide.shapes.add_shape(1, 0, 0, Inches(2.2), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(37, 99, 235)
    bar.line.fill.background()
    
    box = slide.shapes.add_textbox(Inches(0.3), Inches(2.5), Inches(1.6), Inches(3))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "核心概念"
    set_font(p.runs[0], 32, True, (255,255,255))
    
    box = slide.shapes.add_textbox(Inches(2.6), Inches(0.5), Inches(10), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_font(p.runs[0], 36, True)
    
    box = slide.shapes.add_textbox(Inches(2.6), Inches(1.5), Inches(10), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = subtitle
    set_font(p.runs[0], 28, True, (37,99,235))
    
    box = slide.shapes.add_textbox(Inches(2.6), Inches(2.5), Inches(10), Inches(4.5))
    tf = box.text_frame
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        set_font(p.runs[0], 22)
        p.space_before = Pt(20)

def add_flow_slide(slide, title, steps):
    """版面7: 流程步驟"""
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(254, 242, 242)
    bg.line.fill.background()
    
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_font(p.runs[0], 36, True, (220,38,38))
    
    for i, (num, step_title, desc) in enumerate(steps):
        x = Inches(0.8 + i * 4.2)
        
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x+Inches(1.3), Inches(1.8), Inches(1.4), Inches(1.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(220, 38, 38)
        circle.line.fill.background()
        
        box = slide.shapes.add_textbox(x+Inches(1.3), Inches(2.1), Inches(1.4), Inches(0.8))
        p = box.text_frame.paragraphs[0]
        p.text = num
        set_font(p.runs[0], 36, True, (255,255,255))
        p.alignment = PP_ALIGN.CENTER
        
        box = slide.shapes.add_textbox(x, Inches(3.4), Inches(4), Inches(0.7))
        p = box.text_frame.paragraphs[0]
        p.text = step_title
        set_font(p.runs[0], 28, True, (220,38,38))
        p.alignment = PP_ALIGN.CENTER
        
        box = slide.shapes.add_textbox(x, Inches(4.1), Inches(4), Inches(0.8))
        p = box.text_frame.paragraphs[0]
        p.text = desc
        set_font(p.runs[0], 20)
        p.alignment = PP_ALIGN.CENTER
        
        if i < 2:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x+Inches(3.9), Inches(2.3), Inches(0.7), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(200,200,200)
            arrow.line.fill.background()

def add_icon_list_slide(slide, title, items):
    """版面8: 圖標列表"""
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255,255,255)
    bg.line.fill.background()
    
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_font(p.runs[0], 36, True)
    
    colors = [(59,130,246), (111,66,193), (220,38,38), (217,119,6), (25,135,84)]
    for i, (emoji, item_title, desc) in enumerate(items):
        y = Inches(1.4 + i * 1.15)
        color = colors[i % len(colors)]
        
        bar = slide.shapes.add_shape(1, Inches(0.6), y, Pt(6), Inches(1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*color)
        bar.line.fill.background()
        
        box = slide.shapes.add_textbox(Inches(0.9), y+Inches(0.05), Inches(0.8), Inches(0.8))
        p = box.text_frame.paragraphs[0]
        p.text = emoji
        set_font(p.runs[0], 28, False, color)
        
        box = slide.shapes.add_textbox(Inches(1.7), y+Inches(0.05), Inches(4), Inches(0.6))
        p = box.text_frame.paragraphs[0]
        p.text = item_title
        set_font(p.runs[0], 24, True)
        
        box = slide.shapes.add_textbox(Inches(1.7), y+Inches(0.55), Inches(10), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        p.text = desc
        set_font(p.runs[0], 18, False, (100,100,100))

# 主函數
def create_ppt():
    print("🚀 創建多樣化版面 FYP PPT...")
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 頁1: 標題 (版面1)
    add_title_slide(prs.slides.add_slide(prs.slide_layouts[6]), 
        "AI 使用形態的演進", "從 GPT 到 Agent、MCP 與 Skills")
    
    # 頁2: 目錄 (版面8)
    add_icon_list_slide(prs.slides.add_slide(prs.slide_layouts[6]), "目錄", [
        ("📖", "緒論", "研究背景與動機"),
        ("🤖", "GPT 時代", "GPT-3/4 與 ChatGPT"),
        ("🔧", "Agent 革命", "ReAct 與多 Agent 系統"),
        ("⚡", "MCP 與 Skills", "協議與模組化能力"),
        ("🚀", "未來展望", "技術演進與產業影響")
    ])
    
    # 頁3: 研究背景 (版面2)
    add_two_column_slide(prs.slides.add_slide(prs.slide_layouts[6]),
        "🔬", "研究背景", [
            "AI 技術從簡單聊天機器人快速演進到複雜智能代理系統",
            "傳統 GPT 模型擅長文字生成，但缺乏實際行動能力",
            "AI Agent 的出現橋接了理解與執行之間的鴻溝",
            "MCP 和 Skills 系統提供模組化、可擴展的 AI 能力"
        ])
    
    # 頁4: GPT-3 (版面6)
    add_sidebar_slide(prs.slides.add_slide(prs.slide_layouts[6]),
        "GPT-3 時代", "2020 年的突破", [
            "1,750 億參數 — 當時前所未有的規模",
            "少樣本學習能力浮現，展現強大泛化能力",
            "文字生成品質接近人類水平，開啟新紀元",
            "局限性：上下文窗口有限、產生幻覺"
        ])
    
    # 頁5: GPT-4 (版面4)
    add_big_circle_slide(prs.slides.add_slide(prs.slide_layouts[6]),
        "🧠", "GPT-4 突破 (2023)", [
            "多模態能力 — 同時理解文字與圖像",
            "推理能力提升，幻覺問題顯著減少",
            "上下文窗口擴展（8K → 32K → 128K）",
            "更好的指令遵循與安全對齊"
        ])
    
    # 頁6: GPT 局限 (版面3)
    add_three_cards_slide(prs.slides.add_slide(prs.slide_layouts[6]),
        "GPT 模型的局限性", [
            ("無法執行", "只能生成文字，無法執行實際行動"),
            ("無持久記憶", "跨對話無記憶，每次從頭開始"),
            ("推理受限", "複雜任務的推理能力有限")
        ])
    
    # 頁7: Agent 概念 (版面5)
    add_grid_slide(prs.slides.add_slide(prs.slide_layouts[6]),
        "Agent vs 傳統聊天機器人", [
            ("💬", "聊天機器人", "單輪、無狀態、文字輸出"),
            ("🔧", "AI Agent", "多步驟、有狀態、可執行"),
            ("😶", "被動回應", "等待指令，僅回應問題"),
            ("🎯", "主動解決", "理解目標，主動執行")
        ])
    
    # 頁8: ReAct (版面7)
    add_flow_slide(prs.slides.add_slide(prs.slide_layouts[6]),
        "ReAct 框架", [
            ("1", "觀察", "接收環境信息"),
            ("2", "思考", "推理分析狀況"),
            ("3", "行動", "執行具體操作")
        ])
    
    # 頁9: MCP (版面2)
    add_two_column_slide(prs.slides.add_slide(prs.slide_layouts[6]),
        "⚡", "Model Context Protocol", [
            "AI 模型整合的標準化協議",
            "由 Anthropic 於 2024 年提出",
            "實現安全、雙向的通訊機制",
            "好比 AI 應用的 USB-C 接口"
        ])
    
    # 頁10: 核心發現 (版面6)
    add_sidebar_slide(prs.slides.add_slide(prs.slide_layouts[6]),
        "核心發現", "三個關鍵轉變", [
            "1. 演進軌跡：生成 → 推理 → 行動",
            "2. Agent 橋接了 AI 與現實任務的鴻溝",
            "3. MCP 實現標準化、可互操作的 AI 工具",
            "4. Skills 讓 AI 能力模組化與可重用",
            "5. 未來是自主、協作的 AI 系統"
        ])
    
    # 頁11: 致謝 (版面1)
    add_title_slide(prs.slides.add_slide(prs.slide_layouts[6]),
        "致謝", "AI 的未來是 Agentic")
    
    # 保存
    output = 'AI_Evolution_FYP_Varied_v2.pptx'
    prs.save(output)
    
    print(f"✅ 完成！")
    print(f"   總頁數: {len(prs.slides)}")
    print(f"   版面種類: 8 種")
    print(f"   輸出: {output}")
    
    return output

if __name__ == "__main__":
    os.chdir('/Users/singit/.openclaw/workspace-rem/projects/swarm_examples')
    create_ppt()
